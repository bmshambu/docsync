"""Text extraction + chunking.

Ported from the Cowork skill ``skills/rfp-data-prep/extract_text.py`` and made
importable: output directories and chunk settings are passed in rather than
hard-coded relative to a repo root. Logic (paragraph extraction per format,
word-overlap chunking with page/section refs) is unchanged.
"""

from __future__ import annotations

import json
import re
from pathlib import Path


def clean(text: str) -> str:
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def make_chunks(
    paragraphs: list[dict],
    doc_id: str,
    filename: str,
    chunk_size: int,
    overlap: int,
) -> list[dict]:
    """Word-level overlapping chunks carrying page + section refs for citation."""
    chunks: list[dict] = []
    chunk_id = 0
    buffer_words: list[str] = []
    buffer_meta: list[tuple[int, int | None, str]] = []

    for para in paragraphs:
        for w in para["text"].split():
            buffer_meta.append((len(buffer_words), para["page"], para["section"]))
            buffer_words.append(w)

        while len(buffer_words) >= chunk_size:
            chunk_words = buffer_words[:chunk_size]
            meta_slice = buffer_meta[:chunk_size]
            chunks.append(
                {
                    "chunk_id": f"{doc_id}_chunk_{chunk_id:04d}",
                    "doc_id": doc_id,
                    "filename": filename,
                    "page_start": meta_slice[0][1],
                    "page_end": meta_slice[-1][1],
                    "section": meta_slice[0][2],
                    "text": " ".join(chunk_words),
                }
            )
            chunk_id += 1
            buffer_words = buffer_words[chunk_size - overlap :]
            buffer_meta = buffer_meta[chunk_size - overlap :]

    if buffer_words:
        chunks.append(
            {
                "chunk_id": f"{doc_id}_chunk_{chunk_id:04d}",
                "doc_id": doc_id,
                "filename": filename,
                "page_start": buffer_meta[0][1] if buffer_meta else None,
                "page_end": buffer_meta[-1][1] if buffer_meta else None,
                "section": buffer_meta[0][2] if buffer_meta else "",
                "text": " ".join(buffer_words),
            }
        )

    return chunks


# ── Per-format extractors ─────────────────────────────────────────────────────

def extract_docx(path: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(path))
    paragraphs: list[dict] = []
    current_section = "Document Start"
    word_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            current_section = text
        word_count += len(text.split())
        page_estimate = max(1, word_count // 300 + 1)
        paragraphs.append(
            {"text": text, "page": page_estimate, "section": current_section}
        )
    return paragraphs


def extract_pdf(path: Path) -> list[dict]:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    paragraphs: list[dict] = []
    current_section = "Document Start"

    for page_num, page in enumerate(doc, start=1):
        for block in page.get_text("blocks"):
            text = block[4].strip()
            if not text:
                continue
            if len(text) < 80 and text.isupper():
                current_section = text
            paragraphs.append(
                {"text": clean(text), "page": page_num, "section": current_section}
            )
    return paragraphs


def extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    paragraphs: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        current_section = f"Slide {slide_num}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                paragraphs.append(
                    {"text": text, "page": slide_num, "section": current_section}
                )
    return paragraphs


_EXTRACTORS = {
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
}


# ── Public API ────────────────────────────────────────────────────────────────

def process_file(
    path: Path,
    text_dir: Path,
    chunks_dir: Path,
    chunk_size: int,
    overlap: int,
    skip_existing: bool = True,
) -> dict | None:
    """Extract one file → write ``<doc>.txt`` and ``<doc>_chunks.json``.

    Returns a small result dict, or ``None`` for unsupported formats.
    """
    ext = path.suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return None

    doc_id = path.stem.replace(" ", "_")
    text_out = text_dir / f"{doc_id}.txt"
    chunks_out = chunks_dir / f"{doc_id}_chunks.json"

    if skip_existing and text_out.exists() and chunks_out.exists():
        return {
            "doc_id": doc_id,
            "filename": path.name,
            "skipped": True,
            "chunks": len(json.loads(chunks_out.read_text(encoding="utf-8"))),
        }

    paragraphs = extractor(path)
    full_text = "\n\n".join(p["text"] for p in paragraphs)
    text_out.write_text(full_text, encoding="utf-8")

    chunks = make_chunks(paragraphs, doc_id, path.name, chunk_size, overlap)
    chunks_out.write_text(
        json.dumps(chunks, indent=2, ensure_ascii=False), encoding="utf-8"
    )

    return {
        "doc_id": doc_id,
        "filename": path.name,
        "skipped": False,
        "paragraphs": len(paragraphs),
        "chunks": len(chunks),
    }


def extract_all(
    doc_paths: list[Path],
    text_dir: Path,
    chunks_dir: Path,
    chunk_size: int,
    overlap: int,
    skip_existing: bool = True,
    on_progress=None,
) -> list[dict]:
    """Extract a batch of files. ``on_progress(done, total, result)`` is called
    after each file so a job runner can stream progress."""
    text_dir.mkdir(parents=True, exist_ok=True)
    chunks_dir.mkdir(parents=True, exist_ok=True)

    results: list[dict] = []
    total = len(doc_paths)
    for i, path in enumerate(doc_paths, start=1):
        try:
            r = process_file(path, text_dir, chunks_dir, chunk_size, overlap, skip_existing)
            if r is None:
                r = {"filename": path.name, "error": f"unsupported format {path.suffix}"}
        except Exception as exc:  # noqa: BLE001 — log and continue per skill spec
            r = {"filename": path.name, "error": str(exc)}
        results.append(r)
        if on_progress:
            on_progress(i, total, r)
    return results
